#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for Smart Data Processing Assistant
This script creates a professional PowerPoint presentation from the markdown content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def create_presentation():
    """Create the PowerPoint presentation with all slides."""

    # Create presentation object
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    title_content_layout = prs.slide_layouts[1]  # Title and content
    section_header_layout = prs.slide_layouts[2]  # Section header
    two_content_layout = prs.slide_layouts[3]  # Two content
    comparison_layout = prs.slide_layouts[4]  # Comparison
    title_only_layout = prs.slide_layouts[5]  # Title only
    blank_layout = prs.slide_layouts[6]  # Blank
    content_caption_layout = prs.slide_layouts[7]  # Content with caption
    picture_caption_layout = prs.slide_layouts[8]  # Picture with caption

    # Slide content data
    slides_data = [
        # Slide 1: Title Slide
        {
            "layout": title_slide_layout,
            "title": "Smart Data Processing Assistant",
            "subtitle": "Transforming Messy Data into Clean, Searchable Insights",
            "presenter": "Presented by: AI Assistant",
            "date": "Date: April 25, 2026"
        },

        # Slide 2: Agenda
        {
            "layout": title_content_layout,
            "title": "Agenda",
            "content": [
                "• Project Overview & Problem Statement",
                "• Key Features & Capabilities",
                "• Technical Architecture",
                "• Supported File Formats",
                "• Data Processing Pipeline",
                "• User Interface & Experience",
                "• Installation & Setup",
                "• Demo & Use Cases",
                "• Future Enhancements",
                "• Q&A"
            ]
        },

        # Slide 3: Problem Statement
        {
            "layout": title_content_layout,
            "title": "Problem Statement",
            "content": [
                "The Challenge:",
                "• Organizations deal with messy, inconsistent data from multiple sources",
                "• Manual data cleaning is time-consuming and error-prone",
                "• Different file formats require different processing approaches",
                "• Users need quick access to specific information within large datasets",
                "• Exporting and sharing cleaned data is cumbersome",
                "",
                "Our Solution: An intelligent, automated data processing assistant"
            ]
        },

        # Slide 4: Project Overview
        {
            "layout": title_content_layout,
            "title": "Project Overview",
            "content": [
                "Smart Data Processing Assistant",
                "A comprehensive solution for automated data cleaning, processing, and analysis",
                "",
                "Core Capabilities:",
                "✅ Multi-format file processing (CSV, Excel, JSON, TXT, PDF)",
                "✅ AI-powered data cleaning and standardization",
                "✅ Intelligent search and filtering",
                "✅ Multiple export formats",
                "✅ Email delivery system",
                "✅ User-friendly web interface"
            ]
        },

        # Slide 5: Key Features
        {
            "layout": title_content_layout,
            "title": "Key Features Overview",
            "content": [
                "🔄 Data Cleaning & Standardization",
                "• Automatic file type detection",
                "• Text normalization and cleaning",
                "• Column name standardization (snake_case)",
                "• Data type inference and conversion",
                "• Missing value handling",
                "• Duplicate detection and removal",
                "",
                "🔍 Smart Search & Filtering",
                "• Exact matching for IDs and account numbers",
                "• AI-powered semantic search using embeddings",
                "• Natural language query processing"
            ]
        },

        # Slide 6: More Features
        {
            "layout": title_content_layout,
            "title": "Key Features (Continued)",
            "content": [
                "📤 Flexible Export Options",
                "• CSV, Excel (.xlsx), JSON, TXT formats",
                "• Filtered result exports",
                "• Proper formatting and readability",
                "",
                "📧 Email Delivery System",
                "• Automated file attachments",
                "• Customizable email templates",
                "• SMTP configuration support",
                "",
                "🛡️ Error Handling & Validation",
                "• Comprehensive error messages",
                "• File validation and type checking",
                "• Robust processing for edge cases"
            ]
        },

        # Slide 7: Supported Formats
        {
            "layout": two_content_layout,
            "title": "Supported File Formats",
            "content_left": [
                "📁 Input Formats:",
                "• CSV - Comma-separated values",
                "• Excel - .xlsx and .xls files",
                "• JSON - Structured data objects",
                "• TXT - Plain text with auto-parsing",
                "• PDF - Text extraction from documents"
            ],
            "content_right": [
                "📤 Output Formats:",
                "• CSV - Standard format",
                "• Excel - Native .xlsx with formatting",
                "• JSON - Structured records array",
                "• TXT - Tab-separated format"
            ]
        },

        # Slide 8: Technical Architecture
        {
            "layout": title_content_layout,
            "title": "Technical Architecture",
            "content": [
                "🏗️ System Architecture:",
                "",
                "┌─────────────────┐    ┌──────────────────┐    ┌─────────────────┐",
                "│   File Upload   │───▶│  Data Cleaning   │───▶│   Search Index  │",
                "│   (Streamlit)   │    │   Pipeline       │    │   (FAISS)       │",
                "└─────────────────┘    └──────────────────┘    └─────────────────┘",
                "         │                       │                       │",
                "         ▼                       ▼                       ▼",
                "┌─────────────────┐    ┌──────────────────┐    ┌─────────────────┐",
                "│   Export        │    │  Email Delivery  │    │   Results       │",
                "│   System        │    │   (SMTP)         │    │   Display       │",
                "└─────────────────┘    └──────────────────┘    └─────────────────┘"
            ]
        },

        # Slide 9: Data Processing Pipeline
        {
            "layout": title_content_layout,
            "title": "Data Processing Pipeline",
            "content": [
                "🔄 Complete Processing Workflow:",
                "",
                "1. File Ingestion",
                "   • Automatic format detection",
                "   • Schema inference",
                "   • Initial data loading",
                "",
                "2. Data Cleaning",
                "   • Text normalization",
                "   • Column standardization",
                "   • Data type conversion",
                "   • Missing value imputation",
                "",
                "3. Validation & Quality Checks",
                "   • Business rule validation",
                "   • Duplicate detection",
                "   • Data quality scoring"
            ]
        },

        # Slide 10: AI-Powered Search
        {
            "layout": title_content_layout,
            "title": "AI-Powered Search",
            "content": [
                "🧠 Intelligent Search Capabilities:",
                "",
                "Exact Matching:",
                "• Property IDs (R12345 format)",
                "• Account numbers",
                "• Key field lookups",
                "",
                "Semantic Search:",
                "• Natural language queries",
                "• AI-powered similarity matching",
                "• Sentence transformer embeddings",
                "• FAISS vector database",
                "",
                "Query Examples:",
                "• 'Find properties in Dallas'",
                "• 'Show account R12345'",
                "• 'Properties with high tax values'"
            ]
        },

        # Slide 11: User Interface
        {
            "layout": title_content_layout,
            "title": "User Interface",
            "content": [
                "💻 Streamlit Web Interface:",
                "",
                "Main Features:",
                "• Drag-and-drop file upload",
                "• Real-time processing feedback",
                "• Interactive data preview",
                "• Search and filter controls",
                "• Export format selection",
                "• Email delivery options",
                "",
                "User Experience:",
                "• Clean, intuitive design",
                "• Progress indicators",
                "• Error handling with clear messages",
                "• Responsive layout"
            ]
        },

        # Slide 12: Installation & Setup
        {
            "layout": title_content_layout,
            "title": "Installation & Setup",
            "content": [
                "🚀 Getting Started:",
                "",
                "Prerequisites:",
                "• Python 3.10+",
                "• pip package manager",
                "",
                "Installation Steps:",
                "```bash",
                "# Clone repository",
                "git clone https://github.com/username/smart-data-processing-assistant.git",
                "cd smart-data-processing-assistant",
                "",
                "# Install dependencies",
                "pip install -r requirements.txt",
                "",
                "# Run application",
                "streamlit run app.py",
                "```"
            ]
        },

        # Slide 13: Configuration
        {
            "layout": title_content_layout,
            "title": "Configuration",
            "content": [
                "⚙️ Optional Configuration:",
                "",
                "Email Setup (for delivery feature):",
                "```bash",
                "# Environment variables",
                "export EMAIL_USER=your_email@gmail.com",
                "export EMAIL_PASS=your_app_password",
                "```",
                "",
                "Supported Email Providers:",
                "• Gmail (with App Passwords)",
                "• Outlook",
                "• Custom SMTP servers"
            ]
        },

        # Slide 14: Demo - File Upload
        {
            "layout": title_content_layout,
            "title": "Demo - File Upload",
            "content": [
                "📁 File Upload Process:",
                "",
                "1. User selects file (CSV, Excel, JSON, TXT, PDF)",
                "2. System detects file type automatically",
                "3. Shows processing progress",
                "4. Displays cleaned data preview",
                "5. Enables search functionality",
                "",
                "[Screenshot Placeholder]",
                "*[Include screenshot of file upload interface]*"
            ]
        },

        # Slide 15: Demo - Data Preview
        {
            "layout": title_content_layout,
            "title": "Demo - Data Preview",
            "content": [
                "📊 Data Preview Features:",
                "",
                "• Column information table",
                "• Data type detection",
                "• Row count and statistics",
                "• First 20 rows display",
                "• Validation issue reporting",
                "",
                "[Screenshot Placeholder]",
                "*[Include screenshot of data preview]*"
            ]
        },

        # Slide 16: Demo - Search
        {
            "layout": title_content_layout,
            "title": "Demo - Search Functionality",
            "content": [
                "🔍 Search Interface:",
                "",
                "• Text input for queries",
                "• Exact match results",
                "• AI-powered semantic search",
                "• Filtered result display",
                "• Export options for results",
                "",
                "[Screenshot Placeholder]",
                "*[Include screenshot of search results]*"
            ]
        },

        # Slide 17: Use Cases
        {
            "layout": title_content_layout,
            "title": "Use Cases",
            "content": [
                "🎯 Real-World Applications:",
                "",
                "Financial Services:",
                "• Customer data cleaning and validation",
                "• Account reconciliation",
                "• Regulatory reporting",
                "",
                "Real Estate:",
                "• Property data standardization",
                "• Market analysis preparation",
                "• Portfolio management",
                "",
                "Healthcare:",
                "• Patient data processing",
                "• Medical record cleaning",
                "• Research data preparation"
            ]
        },

        # Slide 18: More Use Cases
        {
            "layout": title_content_layout,
            "title": "Use Cases (Continued)",
            "content": [
                "📈 Business Intelligence:",
                "• Sales data consolidation",
                "• Customer segmentation",
                "• Performance analytics",
                "",
                "Research & Academia:",
                "• Survey data processing",
                "• Experimental data cleaning",
                "• Dataset preparation",
                "",
                "Government & Public Sector:",
                "• Public record processing",
                "• Census data validation",
                "• Administrative data management"
            ]
        },

        # Slide 19: Technical Stack
        {
            "layout": title_content_layout,
            "title": "Technical Stack",
            "content": [
                "🛠️ Technology Stack:",
                "",
                "Core Technologies:",
                "• Python 3.10+ - Main programming language",
                "• Streamlit - Web interface framework",
                "• Pandas - Data manipulation",
                "• NumPy - Numerical computing",
                "",
                "AI/ML Components:",
                "• Sentence Transformers - Text embeddings",
                "• FAISS - Vector similarity search",
                "• Scikit-learn - Machine learning utilities"
            ]
        },

        # Slide 20: More Technical Stack
        {
            "layout": title_content_layout,
            "title": "Technical Stack (Continued)",
            "content": [
                "📚 Libraries & Dependencies:",
                "",
                "Data Processing:",
                "• OpenPyXL - Excel file handling",
                "• PyPDF2 - PDF text extraction",
                "• Tabulate - Data formatting",
                "",
                "Development Tools:",
                "• Pytest - Testing framework",
                "• Black - Code formatting",
                "• Pylint - Code quality",
                "",
                "Deployment:",
                "• Docker (future)",
                "• GitHub Actions (future)"
            ]
        },

        # Slide 21: Performance & Scalability
        {
            "layout": title_content_layout,
            "title": "Performance & Scalability",
            "content": [
                "⚡ Performance Characteristics:",
                "",
                "Processing Speed:",
                "• Small files (< 1MB): Instant processing",
                "• Medium files (1-10MB): < 30 seconds",
                "• Large files (10MB+): < 2 minutes",
                "",
                "Search Performance:",
                "• Exact matching: Sub-second response",
                "• Semantic search: < 5 seconds",
                "• Index building: Scales with data size",
                "",
                "Memory Usage:",
                "• Efficient pandas operations",
                "• Streaming for large files",
                "• Configurable batch processing"
            ]
        },

        # Slide 22: Security & Privacy
        {
            "layout": title_content_layout,
            "title": "Security & Privacy",
            "content": [
                "🔒 Security Features:",
                "",
                "Data Handling:",
                "• Local processing (no external uploads)",
                "• Temporary file cleanup",
                "• No persistent data storage",
                "",
                "Email Security:",
                "• SMTP with TLS encryption",
                "• App password support (OAuth-ready)",
                "• Configurable credentials",
                "",
                "Code Security:",
                "• Input validation and sanitization",
                "• Safe file type checking",
                "• Error handling without information leakage"
            ]
        },

        # Slide 23: Testing & Quality Assurance
        {
            "layout": title_content_layout,
            "title": "Testing & Quality Assurance",
            "content": [
                "🧪 Quality Assurance:",
                "",
                "Test Coverage:",
                "• Unit tests for all modules",
                "• Integration tests for pipelines",
                "• End-to-end testing for UI",
                "",
                "Test Framework:",
                "• Pytest - Test execution",
                "• Coverage.py - Code coverage reporting",
                "• GitHub Actions - CI/CD pipeline",
                "",
                "Code Quality:",
                "• Type hints throughout codebase",
                "• Docstrings for all functions",
                "• PEP 8 compliance"
            ]
        },

        # Slide 24: Future Enhancements
        {
            "layout": title_content_layout,
            "title": "Future Enhancements",
            "content": [
                "🚀 Roadmap:",
                "",
                "Short Term (Next 3 months):",
                "• Docker containerization",
                "• GitHub Actions CI/CD",
                "• Advanced data visualization",
                "• Batch processing capabilities",
                "",
                "Medium Term (3-6 months):",
                "• Multi-language support",
                "• Advanced AI features (GPT integration)",
                "• Cloud storage integration",
                "• API endpoints for automation",
                "",
                "Long Term (6+ months):",
                "• Real-time data streaming",
                "• Advanced analytics dashboard",
                "• Machine learning model training",
                "• Enterprise features (user management, audit logs)"
            ]
        },

        # Slide 25: Contributing & Community
        {
            "layout": title_content_layout,
            "title": "Contributing & Community",
            "content": [
                "🤝 How to Contribute:",
                "",
                "Development Setup:",
                "```bash",
                "# Fork repository",
                "# Clone your fork",
                "git clone https://github.com/yourusername/smart-data-processing-assistant.git",
                "",
                "# Create feature branch",
                "git checkout -b feature/new-feature",
                "",
                "# Make changes and test",
                "# Submit pull request",
                "```",
                "",
                "Contribution Guidelines:",
                "• Follow PEP 8 style guide",
                "• Add tests for new features",
                "• Update documentation",
                "• Use meaningful commit messages"
            ]
        },

        # Slide 26: License & Support
        {
            "layout": title_content_layout,
            "title": "License & Support",
            "content": [
                "📄 Project Information:",
                "",
                "License:",
                "• MIT License (open source)",
                "• Free for commercial and personal use",
                "",
                "Support:",
                "• GitHub Issues for bug reports",
                "• Documentation in repository",
                "• Community discussions",
                "",
                "Repository:",
                "• GitHub: https://github.com/username/smart-data-processing-assistant",
                "• Documentation: README.md",
                "• Issues: GitHub Issues tab"
            ]
        },

        # Slide 27: Summary
        {
            "layout": title_content_layout,
            "title": "Summary",
            "content": [
                "🎯 Key Takeaways:",
                "",
                "✅ Comprehensive Solution: Handles multiple file formats with intelligent processing",
                "",
                "✅ AI-Powered: Semantic search and automated data cleaning",
                "",
                "✅ User-Friendly: Simple web interface with powerful capabilities",
                "",
                "✅ Flexible: Multiple export options and email delivery",
                "",
                "✅ Production-Ready: Robust error handling and testing",
                "",
                "✅ Extensible: Clean architecture for future enhancements"
            ]
        },

        # Slide 28: Q&A
        {
            "layout": title_only_layout,
            "title": "Q&A",
            "content": [
                "Questions & Discussion",
                "",
                "Thank you for your attention!",
                "",
                "Contact Information:",
                "• GitHub: @username",
                "• Repository: smart-data-processing-assistant",
                "• Email: [your-email@example.com]",
                "",
                "Demo Available:",
                "Run `streamlit run app.py` to try it yourself!"
            ]
        },

        # Slide 29: Appendix - Code Examples
        {
            "layout": title_content_layout,
            "title": "Appendix - Code Examples",
            "content": [
                "💻 Sample Usage:",
                "",
                "```python",
                "from property_data_cleaning.processor import PropertyDataCleaner",
                "from property_data_cleaning.file_handler import FileHandler",
                "",
                "# Load and clean data",
                "cleaner = PropertyDataCleaner()",
                "file_handler = FileHandler()",
                "",
                "df, file_type = file_handler.load_file(uploaded_file)",
                "cleaned_df = cleaner.run(df)",
                "",
                "# Export results",
                "data_bytes, filename = file_handler.export_data(cleaned_df, 'csv')",
                "```"
            ]
        },

        # Slide 30: Appendix - API Reference
        {
            "layout": title_content_layout,
            "title": "Appendix - API Reference",
            "content": [
                "🔧 Core Classes:",
                "",
                "PropertyDataCleaner:",
                "• run(df) - Complete cleaning pipeline",
                "• clean(df) - Data cleaning only",
                "• validate(df) - Validation only",
                "",
                "FileHandler:",
                "• load_file(file) - Load various file formats",
                "• export_data(df, format) - Export to specified format",
                "",
                "EmailSender:",
                "• send_email(to, subject, body, attachment) - Send emails with attachments"
            ]
        },

        # Slide 31: Appendix - Configuration
        {
            "layout": title_content_layout,
            "title": "Appendix - Configuration Options",
            "content": [
                "⚙️ Configuration:",
                "",
                "Environment Variables:",
                "```",
                "EMAIL_USER=your_email@gmail.com",
                "EMAIL_PASS=your_app_password",
                "EMAIL_SERVER=smtp.gmail.com",
                "EMAIL_PORT=587",
                "```",
                "",
                "Streamlit Configuration:",
                "```",
                "# In .streamlit/config.toml",
                "[server]",
                "headless = true",
                "port = 8501",
                "```"
            ]
        },

        # Slide 32: Thank You
        {
            "layout": title_only_layout,
            "title": "Thank You!",
            "content": [
                "🎉 Thank You!",
                "",
                "Questions? Comments? Feedback?",
                "",
                "Let's connect and collaborate!",
                "",
                "*End of Presentation*"
            ]
        }
    ]

    # Create slides
    for slide_data in slides_data:
        slide = prs.slides.add_slide(slide_data["layout"])

        # Set title
        if "title" in slide_data:
            title = slide.shapes.title
            title.text = slide_data["title"]

        # Handle different slide layouts
        if slide_data["layout"] == title_slide_layout:
            # Title slide with subtitle
            if "subtitle" in slide_data:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = slide_data["subtitle"]

            # Add presenter and date if available
            if len(slide.placeholders) > 2:
                content_placeholder = slide.placeholders[2]
                content_text = ""
                if "presenter" in slide_data:
                    content_text += slide_data["presenter"] + "\n"
                if "date" in slide_data:
                    content_text += slide_data["date"]
                content_placeholder.text = content_text

        elif slide_data["layout"] == title_content_layout:
            # Title and content layout
            content_placeholder = slide.placeholders[1]
            content_frame = content_placeholder.text_frame
            content_frame.text = ""

            for item in slide_data.get("content", []):
                p = content_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(14)

        elif slide_data["layout"] == two_content_layout:
            # Two content layout
            left_placeholder = slide.placeholders[1]
            right_placeholder = slide.placeholders[2]

            # Left content
            left_frame = left_placeholder.text_frame
            left_frame.text = ""
            for item in slide_data.get("content_left", []):
                p = left_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(12)

            # Right content
            right_frame = right_placeholder.text_frame
            right_frame.text = ""
            for item in slide_data.get("content_right", []):
                p = right_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(12)

        elif slide_data["layout"] == title_only_layout:
            # Title only layout - add content as paragraphs
            if "content" in slide_data and len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
                content_frame = content_placeholder.text_frame
                content_frame.text = ""

                for item in slide_data["content"]:
                    p = content_frame.add_paragraph()
                    p.text = item
                    p.font.size = Pt(18)
                    p.alignment = PP_ALIGN.CENTER

    # Save the presentation
    prs.save("smart_data_processing_assistant_presentation.pptx")
    print("✅ PowerPoint presentation created successfully!")
    print("📁 File saved as: smart_data_processing_assistant_presentation.pptx")

if __name__ == "__main__":
    create_presentation()